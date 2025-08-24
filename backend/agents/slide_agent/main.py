from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
from typing import Dict, Any, Optional
import requests
import tempfile
import os

from a2a.types import AgentCard, AgentSkill
from a2a.server.agent_execution import AgentExecutor
from a2a.server.request_handlers import DefaultRequestHandler
from a2a.server.apps import A2AStarletteApplication

from shared.models import AgentRequest, AgentResponse, SlideContent, SlideAgenda
from shared.storage import blob_client
from shared.config import settings


class SlideCreationExecutor(AgentExecutor):
    def __init__(self):
        pass
    
    async def execute(self, request: AgentRequest) -> AgentResponse:
        """スライド作成を実行"""
        try:
            payload = request.payload
            agenda_data = payload.get("agenda", {})
            information = payload.get("information", {})
            template_id = payload.get("template_id")
            include_images = payload.get("include_images", True)
            include_tables = payload.get("include_tables", True)
            
            agenda = SlideAgenda(**agenda_data)
            
            # スライド作成
            pptx_data = await self._create_presentation(
                agenda, information, template_id, include_images, include_tables
            )
            
            # Blob Storage にアップロード
            filename = f"presentation_{agenda.slides[0].title[:20]}.pptx"
            blob_url = blob_client.upload_bytes(
                pptx_data, filename, request.user_id, "presentations"
            )
            
            return AgentResponse(
                request_id=request.request_id,
                success=True,
                result={"slide_url": blob_url, "filename": filename}
            )
        
        except Exception as e:
            return AgentResponse(
                request_id=request.request_id,
                success=False,
                error=str(e)
            )
    
    async def cancel(self, request_id: str) -> bool:
        """処理をキャンセル"""
        return True
    
    async def _create_presentation(
        self, 
        agenda: SlideAgenda, 
        information: Dict[str, Any],
        template_id: Optional[str],
        include_images: bool,
        include_tables: bool
    ) -> bytes:
        """PowerPoint プレゼンテーションを作成"""
        
        # テンプレートまたは新規プレゼンテーション
        if template_id:
            prs = await self._load_template(template_id)
        else:
            prs = Presentation()
        
        # 既存のスライドをクリア（テンプレートの場合）
        if template_id and len(prs.slides) > 0:
            # 最初のスライドをマスターとして保持
            slide_layouts = prs.slide_layouts
        else:
            slide_layouts = prs.slide_layouts
        
        # スライド作成
        for slide_content in agenda.slides:
            self._create_slide(
                prs, slide_content, information, 
                include_images, include_tables, slide_layouts
            )
        
        # バイナリデータとして出力
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
    
    async def _load_template(self, template_id: str) -> Presentation:
        """テンプレートを読み込み"""
        # Cosmos DB からテンプレート情報を取得
        from ...shared.storage import cosmos_client
        
        # 簡略化のため、直接 Presentation を作成
        # 実際の実装では template_id から Blob URL を取得してダウンロード
        return Presentation()
    
    def _create_slide(
        self, 
        prs: Presentation, 
        slide_content: SlideContent,
        information: Dict[str, Any],
        include_images: bool,
        include_tables: bool,
        slide_layouts
    ):
        """個別スライドを作成"""
        
        # スライドレイアウトを選択
        if slide_content.page_number == 1:
            # タイトルスライド
            slide_layout = slide_layouts[0]  # Title Slide
        else:
            # コンテンツスライド
            slide_layout = slide_layouts[1]  # Title and Content
        
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル設定
        title = slide.shapes.title
        if title:
            title.text = slide_content.title
            self._format_title(title)
        
        # コンテンツ設定
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            if content_placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                content_placeholder.text = slide_content.content
                self._format_content(content_placeholder)
        
        # 詳細情報があれば追加
        detailed_content = information.get(f"slide_{slide_content.page_number}", {})
        if detailed_content:
            self._add_detailed_content(slide, detailed_content, include_images, include_tables)
        
        # ノート追加（ハルシネーション警告等）
        if slide_content.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.notes
    
    def _format_title(self, title_shape):
        """タイトルの書式設定"""
        if title_shape.has_text_frame:
            text_frame = title_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.bold = True
    
    def _format_content(self, content_shape):
        """コンテンツの書式設定"""
        if content_shape.has_text_frame:
            text_frame = content_shape.text_frame
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.space_after = Pt(12)
    
    def _add_detailed_content(
        self, 
        slide, 
        detailed_content: Dict[str, Any],
        include_images: bool,
        include_tables: bool
    ):
        """詳細コンテンツを追加"""
        
        # テキストコンテンツの追加
        text_content = detailed_content.get("text", "")
        if text_content and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            if content_placeholder.text:
                content_placeholder.text += "\n\n" + text_content
            else:
                content_placeholder.text = text_content
        
        # 画像の追加
        if include_images and "images" in detailed_content:
            for i, image_url in enumerate(detailed_content["images"][:2]):  # 最大2つの画像
                self._add_image_to_slide(slide, image_url, i)
        
        # テーブルの追加
        if include_tables and "tables" in detailed_content:
            for table_data in detailed_content["tables"][:1]:  # 最大1つのテーブル
                self._add_table_to_slide(slide, table_data)
    
    def _add_image_to_slide(self, slide, image_url: str, position: int):
        """スライドに画像を追加"""
        try:
            # 画像をダウンロード
            response = requests.get(image_url, timeout=10)
            if response.status_code == 200:
                # 一時ファイルに保存
                with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp_file:
                    tmp_file.write(response.content)
                    tmp_file_path = tmp_file.name
                
                # 画像サイズと位置を計算
                left = Inches(6) if position == 0 else Inches(6)
                top = Inches(2 + position * 2.5)
                width = Inches(3)
                
                # スライドに画像を追加
                slide.shapes.add_picture(tmp_file_path, left, top, width=width)
                
                # 一時ファイルを削除
                os.unlink(tmp_file_path)
        except Exception as e:
            print(f"Failed to add image {image_url}: {e}")
    
    def _add_table_to_slide(self, slide, table_data: Dict[str, Any]):
        """スライドにテーブルを追加"""
        try:
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            
            if not headers or not rows:
                return
            
            # テーブルの位置とサイズ
            left = Inches(0.5)
            top = Inches(4)
            width = Inches(9)
            height = Inches(2)
            
            # テーブル作成
            shape = slide.shapes.add_table(
                len(rows) + 1, len(headers), left, top, width, height
            )
            table = shape.table
            
            # ヘッダー設定
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                cell.text_frame.paragraphs[0].font.bold = True
            
            # データ行設定
            for row_idx, row in enumerate(rows):
                for col_idx, cell_data in enumerate(row):
                    if col_idx < len(headers):
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_data)
        
        except Exception as e:
            print(f"Failed to add table: {e}")


# Agent setup
agent_card = AgentCard(
    name="Slide Creation Agent",
    description="python-pptx を使用してPowerPointスライドを作成するエージェント",
    version="1.0.0",
    url="http://slide-agent:8003",
    skills=[],  # Will be populated below
    capabilities={},
    default_input_modes=[],
    default_output_modes=[]
)

agent_skills = [
    AgentSkill(
        name="create_slides",
        description="アジェンダと情報からPowerPointスライドを作成",
        id="create_slides",
        tags=[],
        input_modes=[],
        output_modes=[],
        examples=[]
    )
]

executor = SlideCreationExecutor()
request_handler = DefaultRequestHandler(agent_card, agent_skills, executor)

# Create A2A application
a2a_app = A2AStarletteApplication(agent_card, request_handler)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(a2a_app, host="0.0.0.0", port=8003)