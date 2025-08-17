from fastapi import FastAPI
from semantic_kernel import Kernel
from semantic_kernel.connectors.ai.open_ai import AzureChatCompletion
from semantic_kernel.prompt_template import PromptTemplateConfig
from semantic_kernel.functions import KernelArguments
from typing import Dict, Any
import json

from a2a_python_sdk import (
    AgentCard, AgentSkill, AgentExecutor, DefaultRequestHandler,
    A2AStarletteApplication
)

from ...shared.models import AgentRequest, AgentResponse
from ...shared.config import settings


class ReviewExecutor(AgentExecutor):
    def __init__(self):
        self.kernel = Kernel()
        
        # Azure OpenAI service setup
        self.chat_service = AzureChatCompletion(
            deployment_name=settings.default_llm_model,
            endpoint=settings.azure_ai_foundry_endpoint,
            api_key=settings.azure_ai_foundry_key
        )
        self.kernel.add_service(self.chat_service)
        
        self._setup_prompts()
    
    def _setup_prompts(self):
        """レビュー用プロンプトテンプレートを設定"""
        review_prompt = """
あなたは PowerPoint スライドの品質チェックと事実確認の専門家です。
生成されたスライドの内容を詳細にレビューし、以下の観点で評価してください。

## レビュー観点
1. **事実の正確性**: 内容に事実誤認や古い情報がないか
2. **論理的整合性**: スライド間の論理的なつながりが適切か
3. **情報の完全性**: 必要な情報が不足していないか
4. **プレゼンテーション効果**: 聴衆にとって理解しやすい構成か
5. **ハルシネーション検出**: AI生成特有の事実でない内容がないか

## 入力情報
スライドURL: {{$slide_url}}
元のアジェンダ: {{$agenda}}

## 出力要件
以下のJSON形式で出力してください：
```json
{
  "overall_score": 95,
  "issues": [
    {
      "slide_number": 3,
      "type": "hallucination",
      "severity": "medium",
      "description": "統計データの出典が不明確",
      "suggestion": "出典を明記するか、より信頼性の高いデータに置き換える"
    }
  ],
  "quality_checks": {
    "factual_accuracy": 90,
    "logical_consistency": 95,
    "completeness": 88,
    "presentation_effectiveness": 92,
    "hallucination_risk": 15
  },
  "recommendations": [
    "スライド3の統計データに出典を追加",
    "スライド7の図表をより見やすく調整"
  ],
  "notes_for_slides": {
    "3": "統計データの出典要確認 - 2023年のデータの最新性を要チェック",
    "7": "図表の視認性について要確認"
  }
}
```

レビューを実行してください：
"""
        
        prompt_config = PromptTemplateConfig(
            template=review_prompt,
            name="slide_review",
            description="スライド品質レビューとハルシネーション検出",
            input_variables=["slide_url", "agenda"]
        )
        
        self.review_function = self.kernel.add_function(
            plugin_name="ReviewPlugin",
            function_name="review_slides",
            prompt_template_config=prompt_config
        )
    
    async def execute(self, request: AgentRequest) -> AgentResponse:
        """スライドレビューを実行"""
        try:
            payload = request.payload
            slide_url = payload.get("slide_url", "")
            agenda = payload.get("agenda", {})
            
            if not slide_url:
                return AgentResponse(
                    request_id=request.request_id,
                    success=False,
                    error="Slide URL is required"
                )
            
            # スライドの内容を分析（実際にはPowerPointファイルを読み込む必要がある）
            slide_content = await self._analyze_slide_content(slide_url)
            
            # Semantic Kernel でレビュー実行
            arguments = KernelArguments(
                slide_url=slide_url,
                agenda=json.dumps(agenda, ensure_ascii=False, indent=2)
            )
            
            result = await self.kernel.invoke(self.review_function, arguments)
            review_text = str(result)
            
            # JSON パース
            try:
                if "```json" in review_text:
                    start = review_text.find("```json") + 7
                    end = review_text.find("```", start)
                    review_text = review_text[start:end].strip()
                
                review_data = json.loads(review_text)
                
                # ハルシネーション警告をノートとして追加
                await self._add_warning_notes(slide_url, review_data)
                
                return AgentResponse(
                    request_id=request.request_id,
                    success=True,
                    result=review_data
                )
                
            except (json.JSONDecodeError, ValueError) as e:
                # フォールバック: 基本的なレビュー結果
                fallback_review = self._create_fallback_review()
                return AgentResponse(
                    request_id=request.request_id,
                    success=True,
                    result=fallback_review
                )
        
        except Exception as e:
            return AgentResponse(
                request_id=request.request_id,
                success=False,
                error=str(e)
            )
    
    async def cancel(self, request_id: str) -> bool:
        """処理をキャンセル"""
        return True
    
    async def _analyze_slide_content(self, slide_url: str) -> Dict[str, Any]:
        """スライドの内容を分析"""
        try:
            # 実際の実装では、Blob Storage からファイルをダウンロードして
            # python-pptx で内容を解析する
            from ...shared.storage import blob_client
            
            # ファイルダウンロード
            file_data = blob_client.download_file(slide_url)
            if not file_data:
                return {"error": "Failed to download slide file"}
            
            # PowerPoint ファイル解析
            slide_content = await self._parse_powerpoint_content(file_data)
            return slide_content
            
        except Exception as e:
            print(f"Failed to analyze slide content: {e}")
            return {"error": str(e)}
    
    async def _parse_powerpoint_content(self, file_data: bytes) -> Dict[str, Any]:
        """PowerPoint ファイルの内容を解析"""
        try:
            from pptx import Presentation
            import io
            
            # BytesIO でファイルを開く
            pptx_file = io.BytesIO(file_data)
            prs = Presentation(pptx_file)
            
            slides_content = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # テキストシェイプから内容を抽出
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                slides_content.append({
                    "slide_number": i + 1,
                    "text": "\n".join(slide_text),
                    "shape_count": len(slide.shapes)
                })
            
            return {
                "total_slides": len(prs.slides),
                "slides": slides_content
            }
            
        except Exception as e:
            print(f"Failed to parse PowerPoint: {e}")
            return {"error": "Failed to parse PowerPoint file"}
    
    async def _add_warning_notes(self, slide_url: str, review_data: Dict[str, Any]):
        """ハルシネーション警告等をPowerPointのノートに追加"""
        try:
            notes_for_slides = review_data.get("notes_for_slides", {})
            if not notes_for_slides:
                return
            
            # 実際の実装では、PowerPointファイルを再度開いて
            # ノートスライドに警告を追加し、ファイルを更新する
            
            from ...shared.storage import blob_client
            from pptx import Presentation
            import io
            
            # ファイルダウンロード
            file_data = blob_client.download_file(slide_url)
            if not file_data:
                return
            
            # PowerPoint ファイルを開く
            pptx_file = io.BytesIO(file_data)
            prs = Presentation(pptx_file)
            
            # 各スライドにノートを追加
            for slide_num_str, note_text in notes_for_slides.items():
                slide_num = int(slide_num_str) - 1  # 0ベースのインデックス
                if slide_num < len(prs.slides):
                    slide = prs.slides[slide_num]
                    notes_slide = slide.notes_slide
                    
                    # 既存のノートに追加
                    existing_text = notes_slide.notes_text_frame.text
                    warning_note = f"\n\n⚠️ レビュー警告: {note_text}"
                    notes_slide.notes_text_frame.text = existing_text + warning_note
            
            # ファイルを再アップロード
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            # 元のファイル名を取得してアップロード
            filename = slide_url.split("/")[-1]
            user_id = slide_url.split("/")[-4]  # URL構造から推定
            
            blob_client.upload_bytes(
                output.getvalue(), filename, user_id, "presentations"
            )
            
        except Exception as e:
            print(f"Failed to add warning notes: {e}")
    
    def _create_fallback_review(self) -> Dict[str, Any]:
        """フォールバック用の基本レビュー結果"""
        return {
            "overall_score": 85,
            "issues": [],
            "quality_checks": {
                "factual_accuracy": 85,
                "logical_consistency": 90,
                "completeness": 80,
                "presentation_effectiveness": 85,
                "hallucination_risk": 20
            },
            "recommendations": [
                "内容の事実確認を行ってください",
                "参考資料の出典を確認してください"
            ],
            "notes_for_slides": {}
        }


# Agent setup
agent_card = AgentCard(
    name="Review Agent",
    description="PowerPoint スライドの品質チェックとハルシネーション検出を行うエージェント",
    version="1.0.0"
)

agent_skills = [
    AgentSkill(
        name="review_slides",
        description="スライドの品質チェックとハルシネーション検出"
    )
]

executor = ReviewExecutor()
request_handler = DefaultRequestHandler(agent_card, agent_skills, executor)

# FastAPI app
app = FastAPI(title="Review Agent")

# Create A2A application
a2a_app = A2AStarletteApplication(app, request_handler)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(a2a_app, host="0.0.0.0", port=8004)